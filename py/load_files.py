import json
import os
import sys
from urllib.parse import urlparse
import aiohttp
from io import BytesIO
import asyncio
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from striprtf.striprtf import rtf_to_text
from odf import text
from odf.opendocument import load  # ODF 处理移动到这里避免重复导入
from pptx import Presentation

# 平台检测
IS_WINDOWS = sys.platform == 'win32'
IS_MAC = sys.platform == 'darwin'

# 动态文件类型配置
BASE_OFFICE_EXTS = ['doc', 'docx', 'pptx', 'xls', 'xlsx', 'pdf', 'rtf', 'odt']
PLATFORM_SPECIFIC_EXTS = {
    'win32': ['ppt'],
    'darwin': ['pages', 'numbers', 'key']
}

FILE_FILTERS = [
    { 
        'name': '办公文档', 
        'extensions': BASE_OFFICE_EXTS + PLATFORM_SPECIFIC_EXTS.get(sys.platform, [])
    },
    { 
        'name': '编程开发', 
        'extensions': [
            'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
            'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss',
            'less', 'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml',
            'yaml', 'sql', 'sh'
        ]
    },
    {
        'name': '数据配置',
        'extensions': ['csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml']
    }
]

office_extensions = {ext for group in FILE_FILTERS if group['name'] == '办公文档' for ext in group['extensions']}

async def handle_url(url):
    """异步处理URL输入"""
    async with aiohttp.ClientSession() as session:
        async with session.get(url,headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}) as response:
            response.raise_for_status()
            content = await response.read()
            path = urlparse(url).path
            ext = os.path.splitext(path)[1].lstrip('.').lower()
            return content, ext

async def handle_local_file(file_path):
    """异步处理本地文件"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    loop = asyncio.get_event_loop()
    content = await loop.run_in_executor(None, _read_file, file_path)
    ext = os.path.splitext(file_path)[1].lstrip('.').lower()
    return content, ext

def _read_file(file_path):
    """同步读取文件内容"""
    with open(file_path, 'rb') as f:
        return f.read()

async def get_content(input_str):
    """获取文件内容和扩展名"""
    if input_str.startswith(('http://', 'https://')):
        return await handle_url(input_str)
    else:
        return await handle_local_file(input_str)

def decode_text(content_bytes):
    """通用文本解码（增加BOM处理）"""
    encodings = ['utf-8-sig', 'utf-16', 'gbk', 'iso-8859-1', 'latin-1']
    for enc in encodings:
        try:
            return content_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return content_bytes.decode('utf-8', errors='replace')

async def handle_office_document(content, ext):
    """异步处理办公文档（带平台检测）"""
    handler = {
        'pdf': handle_pdf,
        'docx': handle_docx,
        'xlsx': handle_excel,
        'xls': handle_excel,
        'rtf': handle_rtf,
        'odt': handle_odt,
        'pptx': handle_pptx,
    }
    
    # Windows平台扩展
    if IS_WINDOWS:
        handler['ppt'] = handle_ppt
    
    handler_func = handler.get(ext)
    
    if handler_func:
        return await handler_func(content)
    
    # Mac平台iWork格式处理
    if IS_MAC and ext in ['pages', 'numbers', 'key']:
        raise NotImplementedError(f"iWork格式暂不支持自动解析，请手动导出为通用格式")
    
    raise NotImplementedError(f"暂不支持处理 {ext.upper()} 格式文件")

async def handle_odt(content):
    """异步处理ODT文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_odt, content)

def _process_odt(content):
    """同步处理ODT内容"""
    from odf.teletype import extractText
    
    try:
        doc = load(BytesIO(content))
        text_content = []
        for para in doc.getElementsByType(text.P):
            text_content.append(extractText(para))
        for table in doc.getElementsByType(text.Table):
            for row in table.getElementsByType(text.TableRow):
                row_data = []
                for cell in row.getElementsByType(text.TableCell):
                    row_data.append(extractText(cell))
                text_content.append("\t".join(row_data))
        return '\n'.join(text_content)
    except Exception as e:
        raise RuntimeError(f"ODT文件解析失败: {str(e)}")

async def handle_pdf(content):
    """异步处理PDF文件（增加容错处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pdf, content)

def _process_pdf(content):
    """同步处理PDF内容"""
    text = []
    try:
        with BytesIO(content) as pdf_file:
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                page_text = page.extract_text() or ""  # 处理无文本页面
                text.append(page_text)
    except Exception as e:
        raise RuntimeError(f"PDF解析失败: {str(e)}")
    return '\n'.join(text)

async def handle_docx(content):
    """异步处理DOCX文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_docx, content)

def _process_docx(content):
    """同步处理DOCX内容（增加表格处理）"""
    doc = Document(BytesIO(content))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            text.append('\t'.join(cell.text for cell in row.cells))
    return '\n'.join(text)

async def handle_excel(content):
    """异步处理Excel文件（优化大文件处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_excel, content)

def _process_excel(content):
    """同步处理Excel内容"""
    try:
        wb = load_workbook(filename=BytesIO(content), read_only=True, data_only=True)
        text = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text.append('\t'.join(str(cell) if cell is not None else '' for cell in row))
        return '\n'.join(text)
    except Exception as e:
        raise RuntimeError(f"Excel解析失败: {str(e)}")

async def handle_rtf(content):
    """异步处理RTF文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_rtf, content)

def _process_rtf(content):
    """同步处理RTF内容"""
    try:
        return rtf_to_text(content.decode('utf-8', errors='replace'))
    except Exception as e:
        raise RuntimeError(f"RTF解析失败: {str(e)}")

async def handle_pptx(content):
    """异步处理PPTX文件（优化内容提取）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, content)

def _process_pptx(content):
    """同步处理PPTX内容"""
    try:
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.strip() for cell in row.cells]
                        text.append("\t".join(row_data))
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPTX解析失败: {str(e)}")

async def handle_ppt(content):
    """处理PPT文件（Windows平台专用）"""
    if not IS_WINDOWS:
        raise NotImplementedError("PPT格式仅支持在Windows系统处理")
    
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("请安装pywin32依赖: pip install pywin32")
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_ppt, content)

def _process_ppt(content):
    """同步处理PPT内容（Windows COM API）"""
    import win32com.client
    import tempfile
    import pythoncom

    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(tmp_path)
        text = []
        for slide in pres.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text.append(shape.TextFrame.TextRange.Text.strip())
        pres.Close()
        powerpoint.Quit()
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPT解析失败: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        os.unlink(tmp_path)

async def get_file_content(file_url):
    """异步获取文件内容（增加编码异常处理）"""
    try:
        content, ext = await get_content(file_url)
        if ext in office_extensions:
            return await handle_office_document(content, ext)
        return decode_text(content)
    except Exception as e:
        return f"文件解析错误: {str(e)}"

async def get_files_content(files_path_list):
    """异步获取所有文件内容并拼接（增加错误隔离）"""
    tasks = [get_file_content(fp) for fp in files_path_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for fp, content in zip(files_path_list, contents):
        if isinstance(content, Exception):
            results.append(f"文件 {fp} 解析失败: {str(content)}")
        else:
            results.append(f"文件 {fp} 内容：\n{content}")
    return "\n\n".join(results)

async def get_files_json(files_list):
    """异步获取所有文件内容并拼接为JSON格式（增加错误隔离）
    输入
    files_list: [{'path': 'path/to/file', 'name': 'file_name'}]
    """
    tasks = [get_file_content(files["path"]) for files in files_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for files, content in zip(files_list, contents):
        results.append({"file_path": files["path"],"file_name": files["name"], "content": str(content)})
    return results

ALLOWED_EXTENSIONS = [
  # 办公文档
  'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'pdf', 'pages', 
  'numbers', 'key', 'rtf', 'odt',
  
  # 编程开发
  'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
  'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss', 'less',
  'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml', 'yaml', 
  'sql', 'sh',
  
  # 数据配置
  'csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml'
]

ALLOWED_IMAGE_EXTENSIONS = ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp']

file_tool = {
    "type": "function",
    "function": {
        "name": "get_file_content",
        "description": f"获取给定的文件URL中的内容，无论是公网URL还是服务器内部URL，支持格式：{', '.join(ALLOWED_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "file_url": {
                    "type": "string",
                    "description": "文件URL",
                }
            },
            "required": ["file_url"],
        },
    },
}

image_tool = {
    "type": "function",
    "function": {
        "name": "get_image_content",
        "description": f"获取给定的图片URL中的内容，无论是公网URL还是服务器内部URL，支持格式：{', '.join(ALLOWED_IMAGE_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "image_url": {
                    "type": "string",
                    "description": "图片URL",
                }
            },
            "required": ["image_url"],
        },
    },
}